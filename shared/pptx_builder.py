"""
Shared PPTX builder. Palette, shape helpers, and a presentation factory so
every agent's slide deck looks the same.
"""
from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Brand palette ────────────────────────────────────────────────────────────
VF_GREEN   = RGBColor(0x21, 0x57, 0x32)
VF_OLIVE   = RGBColor(0x65, 0x8D, 0x1B)
VF_GOLD    = RGBColor(0xEA, 0xAA, 0x00)
VF_BURG    = RGBColor(0x6F, 0x26, 0x3D)
VF_BROWN   = RGBColor(0x3E, 0x2B, 0x2E)
VF_TAN     = RGBColor(0xB7, 0xB0, 0x9C)
VF_LTGRAY  = RGBColor(0xE5, 0xE1, 0xE6)
VF_BLUEGR  = RGBColor(0xCE, 0xD9, 0xE5)
VF_DARK    = RGBColor(0x23, 0x1F, 0x20)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
RED_ALERT  = RGBColor(0xC0, 0x39, 0x2B)
AMBER      = RGBColor(0xE8, 0x7C, 0x38)
STABLE_GRN = RGBColor(0x4E, 0xA8, 0x6A)
FONT       = "Arial"

# Standard slide dimensions (widescreen 16:9)
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ── Factory ──────────────────────────────────────────────────────────────────

def create_presentation() -> Presentation:
    """Create a new branded Presentation at standard 16:9 dimensions."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_blank_slide(prs: Presentation):
    """Add a blank slide with white background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    return slide


# ── Shape helpers ────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    """Set a solid background colour on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    """Add a solid-colour rectangle. Returns the shape."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h)
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def hline(slide, l, t, w, color=VF_GREEN):
    """Draw a thin horizontal rule."""
    rect(slide, l, t, w, 0.025, color)


def header_bar(slide, title):
    """Standard slide header: green bar + gold underline + title."""
    rect(slide, 0, 0, 13.333, 0.75, VF_GREEN)
    txt(slide, 0.8, 0.15, 10, 0.45, title, sz=24, color=WHITE, bold=True)
    rect(slide, 0, 0.75, 13.333, 0.04, VF_GOLD)


def title_slide(slide, title, subtitle="", tag=""):
    """Full-height title header (green) with optional subtitle and tag line."""
    rect(slide, 0, 0, 13.333, 1.15, VF_GREEN)
    txt(slide, 0.8, 0.15, 10, 0.55, title, sz=30, color=WHITE, bold=True)
    if subtitle:
        txt(slide, 0.8, 0.65, 8, 0.35, subtitle, sz=13, color=VF_LTGRAY)
    rect(slide, 0, 1.15, 13.333, 0.04, VF_GOLD)
    if tag:
        txt(slide, 0.8, 1.3, 10, 0.3, tag, sz=8, color=VF_TAN)


def txt(slide, l, t, w, h, text, sz=12, color=VF_DARK, bold=False,
        align=PP_ALIGN.LEFT):
    """Add a textbox with a single paragraph. Returns the text_frame."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = FONT
    p.alignment = align
    return tb.text_frame


def bullets(slide, l, t, w, h, items, sz=10.5, color=VF_DARK, bc=VF_GREEN):
    """Add a bulleted list (▸ prefix). Returns the text_frame."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        p.space_before = Pt(1)
        rb = p.add_run()
        rb.text = "\u25B8 "
        rb.font.size = Pt(sz)
        rb.font.color.rgb = bc
        rb.font.name = FONT
        rt = p.add_run()
        rt.text = item
        rt.font.size = Pt(sz)
        rt.font.color.rgb = color
        rt.font.name = FONT
    return tf


def footer(slide, text):
    """Add a small footer line at the bottom of the slide."""
    txt(slide, 0.5, 7.1, 12, 0.25, text, sz=8, color=VF_TAN)


def status_card(slide, x, y, w, h, label, region, detail, color):
    """Add a status card (coloured header stripe + text)."""
    rect(slide, x, y, w, h, VF_GREEN)
    rect(slide, x, y, w, 0.04, color)
    txt(slide, x + 0.1, y + 0.1, w - 0.2, 0.25, label,
        sz=9, color=color, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + 0.1, y + 0.4, w - 0.2, 0.3, region,
        sz=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + 0.1, y + 0.75, w - 0.2, 0.6, detail,
        sz=9, color=VF_LTGRAY, align=PP_ALIGN.CENTER)


def scenario_box(slide, x, y, w, h, title, body, header_color=VF_GREEN):
    """Add a scenario box (Base/Stress/Shock) with coloured header."""
    rect(slide, x, y, w, h, VF_LTGRAY)
    rect(slide, x, y, w, 0.35, header_color)
    txt(slide, x + 0.1, y + 0.05, w - 0.2, 0.25, title,
        sz=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(slide, x + 0.1, y + 0.45, w - 0.2, h - 0.55, body,
        sz=10, color=VF_DARK)
