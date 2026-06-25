"""
Reference slide layouts. Fixed visual patterns for the oil and diesel decks,
built for the 10.0 x 5.62 inch slide format.
"""
from __future__ import annotations
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ── Reference palette ────────────────────────────────────────────────────────
NAVY       = RGBColor(0x0D, 0x1B, 0x2A)  # Dark header/card background
NAVY_LIGHT = RGBColor(0x1A, 0x2F, 0x45)  # Sub-section headers
CARD_BG    = RGBColor(0x0D, 0x1B, 0x2A)  # KPI card background
VF_GREEN   = RGBColor(0x1E, 0x4D, 0x2B)  # diesel header green
TEAL       = RGBColor(0x0E, 0x74, 0x90)  # Accent / positive
RED_CRIT   = RGBColor(0xB9, 0x1C, 0x1C)  # Critical severity
RED_FIRE   = RGBColor(0xC2, 0x41, 0x0C)  # High severity (orange-red)
AMBER_SEV  = RGBColor(0xB4, 0x53, 0x09)  # Amber severity
GOLD       = RGBColor(0xEA, 0xAA, 0x00)  # Gold accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
HINT_GREEN = RGBColor(0xE8, 0xF2, 0xEA)  # Warning banner background
HINT_RED   = RGBColor(0xFD, 0xF0, 0xF0)  # Delta negative background
GREEN_OK   = RGBColor(0xEC, 0xFD, 0xF5)  # Validation OK background
DARK_TEXT  = RGBColor(0x1A, 0x1A, 0x1A)
MUTED_TEXT = RGBColor(0x7A, 0x7A, 0x7A)
FONT       = "Arial"

SW = 10.0   # Slide width in inches
SH = 5.62   # Slide height in inches


def create_ref_presentation() -> Presentation:
    """Create a presentation matching the reference 10x5.62 format."""
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def add_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0x05, 0x0D, 0x15)  # Very dark bg
    return slide


# ── Shape helpers ────────────────────────────────────────────────────────────

def _rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = color; s.line.fill.background()
    return s

def _txt(slide, l, t, w, h, text, sz=12, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.text = text; p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = FONT; p.alignment = align
    return tb.text_frame

def _multi_txt(slide, l, t, w, h, lines, sz=10, color=WHITE, line_spacing=1.1):
    """Add a textbox with multiple paragraphs."""
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color; p.font.name = FONT
        p.space_after = Pt(2)
    return tf


# ── Header & footer (matching reference exactly) ────────────────────────────

def ref_header(slide, title, subtitle="", color=NAVY):
    """0.72" dark header bar with title + subtitle."""
    _rect(slide, 0, 0, SW, 0.72, color)
    _txt(slide, 0.35, 0.08, 7.5, 0.28, title, sz=16, color=WHITE, bold=True)
    if subtitle:
        _txt(slide, 0.35, 0.38, 8.5, 0.24, subtitle, sz=10, color=LIGHT_GRAY)

def ref_footer(slide, left_text, right_text="", color=NAVY):
    """0.33" dark footer bar."""
    _rect(slide, 0, SH - 0.33, SW, 0.33, color)
    _txt(slide, 0.3, SH - 0.31, 7.5, 0.22, left_text, sz=7.5, color=LIGHT_GRAY)
    if right_text:
        _txt(slide, 7.8, SH - 0.31, 2.1, 0.22, right_text, sz=7.5, color=LIGHT_GRAY, align=PP_ALIGN.RIGHT)

def ref_section_bar(slide, y, text, color=NAVY_LIGHT):
    """Thin section header bar."""
    _rect(slide, 0.3, y, SW - 0.6, 0.26, color)
    _txt(slide, 0.5, y + 0.02, SW - 1.0, 0.2, text, sz=9, color=WHITE, bold=True)

def ref_alert_banner(slide, y, text, bg=HINT_GREEN):
    """Light alert banner (e.g. warnings, context notes)."""
    _rect(slide, 0.3, y, SW - 0.6, 0.46, bg)
    _txt(slide, 0.5, y + 0.06, SW - 1.0, 0.35, text, sz=9, color=DARK_TEXT)


# ── KPI card (dark card, large value, label + sub) ───────────────────────────

def kpi_card(slide, x, y, label, value, sublabel, w=2.2, h=1.0, accent=TEAL):
    """Dark KPI card matching VLSFO reference layout."""
    _rect(slide, x, y, w, h, CARD_BG)
    _txt(slide, x + 0.08, y + 0.08, w - 0.16, 0.16, label, sz=8, color=LIGHT_GRAY, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, x + 0.08, y + 0.26, w - 0.16, 0.42, value, sz=26, color=accent, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, x + 0.08, y + 0.72, w - 0.16, 0.16, sublabel, sz=7.5, color=MUTED_TEXT, align=PP_ALIGN.CENTER)


# ── Country price card (diesel reference) ────────────────────────────────────

def country_price_card(slide, x, y, flag_name, price, currency, delta, delta_detail, stripe_color=VF_GREEN):
    """White country card with bottom stripe, matching diesel reference."""
    w, h = 2.2, 1.52
    _rect(slide, x, y, w, h, WHITE)
    _rect(slide, x, y + h - 0.04, w, 0.04, stripe_color)
    _txt(slide, x + 0.14, y + 0.12, w - 0.28, 0.2, flag_name, sz=10, color=DARK_TEXT, bold=True)
    _txt(slide, x + 0.14, y + 0.33, w - 0.28, 0.5, price, sz=28, color=DARK_TEXT, bold=True)
    _txt(slide, x + 0.14, y + 0.84, w - 0.28, 0.18, currency, sz=9, color=MUTED_TEXT)
    _rect(slide, x + 0.14, y + 1.04, w - 0.28, 0.22, HINT_RED)
    _txt(slide, x + 0.14, y + 1.05, w - 0.28, 0.2, delta, sz=9, color=RED_CRIT, bold=True)
    _txt(slide, x + 0.14, y + 1.27, w - 0.28, 0.16, delta_detail, sz=7.5, color=MUTED_TEXT)


# ── Timeline card (horizontal, severity-striped header) ──────────────────────

def timeline_card(slide, x, y, date, text, severity_color=RED_CRIT, w=1.48, h=1.52):
    """Horizontal timeline card matching VLSFO crisis timeline."""
    _rect(slide, x, y, w, h, WHITE)
    _rect(slide, x, y, w, 0.28, severity_color)
    _txt(slide, x + 0.08, y + 0.03, w - 0.16, 0.22, date, sz=9, color=WHITE, bold=True)
    _txt(slide, x + 0.08, y + 0.33, w - 0.16, h - 0.38, text, sz=8, color=DARK_TEXT)


# ── Impact card (emoji + title + body, 4-across) ────────────────────────────

def impact_card(slide, x, y, emoji, title, body, w=2.22, h=1.96):
    """Operational impact card matching VLSFO reference."""
    _rect(slide, x, y, w, h, WHITE)
    _txt(slide, x + 0.12, y + 0.06, 0.34, 0.28, emoji, sz=16, color=DARK_TEXT)
    _txt(slide, x + 0.12, y + 0.36, w - 0.24, 0.26, title, sz=10, color=DARK_TEXT, bold=True)
    _txt(slide, x + 0.12, y + 0.64, w - 0.24, h - 0.74, body, sz=8, color=DARK_TEXT)


# ── Risk register row (severity + title + detail) ───────────────────────────

def risk_row(slide, x, y, severity, title, detail, sev_color=RED_CRIT, w=9.4):
    """Risk register row matching diesel bulletin."""
    _rect(slide, x, y, w, 0.7, WHITE)
    _rect(slide, x, y, w, 0.03, sev_color)
    _txt(slide, x + 0.1, y + 0.08, 1.0, 0.2, severity, sz=8, color=sev_color, bold=True)
    _txt(slide, x + 1.2, y + 0.08, 2.5, 0.2, title, sz=9, color=DARK_TEXT, bold=True)
    _txt(slide, x + 0.1, y + 0.32, w - 0.2, 0.34, detail, sz=8, color=DARK_TEXT)


# ── Scenario box (Base/Stress/Shock) ─────────────────────────────────────────

def scenario_box(slide, x, y, title, body, header_color=VF_GREEN, w=2.9, h=2.2):
    """Scenario box matching VLSFO/diesel forward outlook."""
    _rect(slide, x, y, w, h, WHITE)
    _rect(slide, x, y, w, 0.32, header_color)
    _txt(slide, x + 0.1, y + 0.04, w - 0.2, 0.24, title, sz=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, x + 0.1, y + 0.4, w - 0.2, h - 0.5, body, sz=8, color=DARK_TEXT)


# ── Recovery indicator card ──────────────────────────────────────────────────

def indicator_card(slide, x, y, emoji, title, body, now, target, trigger, w=2.2, h=1.8):
    """Recovery signal card matching VLSFO slide 6."""
    _rect(slide, x, y, w, h, WHITE)
    _txt(slide, x + 0.1, y + 0.06, 0.3, 0.24, emoji, sz=14, color=DARK_TEXT)
    _txt(slide, x + 0.1, y + 0.3, w - 0.2, 0.2, title, sz=9, color=DARK_TEXT, bold=True)
    _txt(slide, x + 0.1, y + 0.52, w - 0.2, 0.6, body, sz=7.5, color=DARK_TEXT)
    # Status lines
    _txt(slide, x + 0.1, y + 1.15, w - 0.2, 0.16, f"NOW: {now}", sz=7, color=TEAL, bold=True)
    _txt(slide, x + 0.1, y + 1.32, w - 0.2, 0.16, f"TARGET: {target}", sz=7, color=MUTED_TEXT)
    _txt(slide, x + 0.1, y + 1.49, w - 0.2, 0.16, f"TRIGGER: {trigger}", sz=7, color=MUTED_TEXT)


# ── Action row (timeline-style with urgency tag) ─────────────────────────────

def action_row(slide, x, y, urgency, text, w=9.0):
    """Action item matching VLSFO/diesel actions slide."""
    urgency_colors = {"IMMEDIATE": RED_CRIT, "THIS WEEK": AMBER_SEV, "30 DAYS": TEAL}
    clr = urgency_colors.get(urgency, MUTED_TEXT)
    _rect(slide, x, y, 1.2, 0.28, clr)
    _txt(slide, x + 0.05, y + 0.03, 1.1, 0.22, urgency, sz=7.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(slide, x + 1.35, y + 0.03, w - 1.4, 0.22, text, sz=8.5, color=WHITE)
